import pdfplumber
from pptx import Presentation
from docx import Document
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import io


# ---------------------- PDF TEXT EXTRACTION ----------------------
def extract_pdf_text(file):
    text_pages = []

    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            extracted = page.extract_text()
            if extracted:
                text_pages.append(extracted.strip())

    return "\n\n".join(text_pages)


# ---------------------- EXPORT DOCX ----------------------
def export_docx(content):
    doc = Document()

    for para in content.split("\n"):
        doc.add_paragraph(para)

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer


# ---------------------- EXPORT PDF ----------------------
def export_pdf(content):
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)

    width, height = letter
    y = height - 50  # Top margin

    for line in content.split("\n"):
        # Auto wrap text for long sentences (manual wrap at 90 chars)
        while len(line) > 90:
            c.drawString(50, y, line[:90])
            line = line[90:]
            y -= 15
            if y < 50:
                c.showPage()
                y = height - 50

        c.drawString(50, y, line)
        y -= 15

        if y < 50:
            c.showPage()
            y = height - 50

    c.save()
    buffer.seek(0)
    return buffer


# ---------------------- EXPORT PPT ----------------------
def export_ppt(ppt_content):
    """
    Expects ppt_content in a structured text format like:
    
    Slide Title
    - Subtopic 1: one line description
    - Subtopic 2: one line description

    Slide Title 2
    - Subtopic A: description
    """

    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    slides = ppt_content.strip().split("\n\n")
    
    for block in slides:
        lines = block.strip().split("\n")
        if not lines:
            continue

        # Extract title
        title_text = lines[0].strip()

        # Slide
        slide_layout = prs.slide_layouts[1]  # Title + content layout
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Body placeholder
        body = slide.shapes.placeholders[1].text_frame
        body.clear()  # Remove default bullet

        # Add bullet points
        for line in lines[1:]:
            if ":" in line:
                topic, desc = line.split(":", 1)
                p = body.add_paragraph()
                p.text = f"{topic.strip()}: {desc.strip()}"
                p.level = 0
                p.font.size = Pt(20)
            else:
                p = body.add_paragraph()
                p.text = line.strip()
                p.level = 0
                p.font.size = Pt(20)

    # Export as binary buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

